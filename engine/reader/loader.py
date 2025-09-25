import chardet
import os
import pytesseract
import shutil
import subprocess
import tempfile

from docx import Document

from engine.reader import asr as asr_processor
from engine.reader import ocr as ocr_instance
from core.logger import log

COMMON_ENCODINGS = ['utf-8', 'gbk', 'gb18030', 'big5', 'latin1']

# 处理文本类型文件
def detect_encoding(file_path):
    """带置信度的编码检测"""
    with open(file_path, 'rb') as f:
        rawdata = f.read(8192)  # 增加样本量

    result = chardet.detect(rawdata)
    return result['encoding'], result['confidence']


def safe_read_text(file_path):
    """多编码尝试的文本读取"""
    # 先尝试检测到的编码
    encoding, confidence = detect_encoding(file_path)

    # 置信度高于70%时优先尝试
    if confidence > 0.7:
        try:
            with open(file_path, 'r', encoding=encoding, errors='strict') as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            pass

    # 降级尝试常见编码
    for enc in COMMON_ENCODINGS:
        try:
            with open(file_path, 'r', encoding=enc, errors='strict') as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue

    # 最终尝试使用替代字符
    try:
        with open(file_path, 'r', encoding=encoding, errors='replace') as f:
            return f.read()
    except:
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='replace')


# 处理doc文件
def read_doc_with_antiword(file_path: str) -> str:
    try:
        result = subprocess.run(
            ['antiword', file_path],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            check=True
        )
        text =  result.stdout.decode('utf-8')
        text += "\n" + ocr_instance.read_office_pic(file_path)
        return text
    except subprocess.CalledProcessError as e:
        raise ValueError(f"读取 .doc 文件失败: {e.stderr.decode()}")


# 将ppt转pptx
def convert_ppt_to_pptx(ppt_path: str) -> str:
    """
    将 PPT 文件转换为 PPTX 格式
    返回临时 PPTX 文件路径
    """
    try:
        # 检查 unoconv 是否可用
        subprocess.run(["unoconv", "--version"], check=True, capture_output=True)
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        raise RuntimeError(
            "请先安装 unoconv 和 LibreOffice:\n"
            "Linux: sudo apt-get install unoconv libreoffice\n"
            "macOS: brew install unoconv libreoffice"
        ) from e

    # 创建临时目录存放转换文件
    temp_dir = tempfile.mkdtemp()
    pptx_path = os.path.join(temp_dir, "converted.pptx")

    try:
        # 执行转换命令
        subprocess.run(
            ["unoconv", "-f", "pptx", "-o", pptx_path, ppt_path],
            check=True,
            stderr=subprocess.PIPE
        )
    except subprocess.CalledProcessError as e:
        shutil.rmtree(temp_dir)  # 清理临时目录
        raise RuntimeError(f"PPT 转换失败: {e.stderr.decode()}") from e

    if not os.path.exists(pptx_path):
        shutil.rmtree(temp_dir)
        raise RuntimeError("转换后的文件未生成")

    return pptx_path


def read_pic(img_path) -> str:
    try:
        return ocr_instance.read(img_path)
    except Exception as e:
        log.error(f"解析 pic 失败 [{img_path}]: {e.__class__.__name__} - {str(e)}")
        return pytesseract.image_to_string(img_path, lang='chi_sim')


def read_pdf(pdf_path) -> str:
    try:
        return ocr_instance.read(pdf_path)
    except Exception as e:
        log.error(f"解析 pdf 失败 [{pdf_path}]: {e.__class__.__name__} - {str(e)}")
        return ""

def extract_text(file_path):
    """支持多种格式的文本提取"""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == '.pdf':
            return read_pdf(file_path)

        elif ext in ('.docx', '.doc'):
            if ext == '.docx':
                text = '\n'.join(p.text for p in Document(file_path).paragraphs)
                text += "\n" + ocr_instance.read_office_pic(file_path)
                return text
            else:
                return read_doc_with_antiword(file_path)

        elif ext in ('.png', '.jpg', '.jpeg', '.bmp'):
            return read_pic(file_path)

        elif ext in ('.txt', '.log', '.csv', '.md', '.json'):
            return safe_read_text(file_path)

        elif ext in ('.pptx', '.ppt'):
            if ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                ret = '\n'.join(text)
                ret += "\n" + ocr_instance.read_office_pic(file_path)
                return ret
            elif ext == '.ppt':
                converted_path = convert_ppt_to_pptx(file_path)  # 调用 unoconv 转换
                return extract_text(converted_path)
        elif ext in ('.mp3', '.wav', '.m4a', '.flac'):
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"音频文件不存在: {file_path}")

            # 执行ASR转录
            if asr_processor.model is None:
                raise RuntimeError("语音识别引擎未初始化")

            return asr_processor.transcribe(file_path)

        else:
            raise ValueError(f"Unsupported format: {ext}")

    except Exception as e:
        log.error(f"解析失败 [{file_path}]: {e.__class__.__name__} - {str(e)}")
        return ""
